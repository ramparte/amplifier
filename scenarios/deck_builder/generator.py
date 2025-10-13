"""PPTX generator for slide decks."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

from .layouts import get_placeholder_indices
from .layouts import select_layout
from .models import Slide
from .models import SlideType
from .models import Theme
from .themes import apply_text_styling
from .themes import get_default_theme
from .themes import hex_to_rgb


def generate_presentation(slides: list[Slide], output_path: str | Path, theme: Theme | None = None) -> None:
    """Generate a PowerPoint presentation from slides.

    Args:
        slides: List of Slide objects
        output_path: Path to save the PPTX file
        theme: Optional theme configuration
    """
    if theme is None:
        theme = get_default_theme()

    # Create presentation
    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Generate each slide
    for slide in slides:
        _add_slide(prs, slide, theme)

    # Save presentation
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _add_slide(presentation: "Presentation", slide: Slide, theme: Theme) -> None:  # type: ignore[name-defined]
    """Add a single slide to the presentation.

    Args:
        presentation: Presentation object
        slide: Slide to add
        theme: Theme configuration
    """
    # Select appropriate layout
    layout_type = select_layout(slide)
    layout_idx = layout_type.value

    # Get slide layout from presentation
    slide_layout = presentation.slide_layouts[layout_idx]
    ppt_slide = presentation.slides.add_slide(slide_layout)

    # Get placeholder indices
    placeholders = get_placeholder_indices(layout_type)

    # Add title
    if "title" in placeholders and slide.title:
        title_placeholder = ppt_slide.placeholders[placeholders["title"]]
        title_placeholder.text = slide.title
        apply_text_styling(title_placeholder.text_frame, theme, is_title=True)

    # Add content based on slide type
    if slide.slide_type == SlideType.CONTENT and "content" in placeholders:
        _add_content(ppt_slide, slide, placeholders["content"], theme)
    elif slide.slide_type == SlideType.TITLE and "subtitle" in placeholders:
        # For title slides, put first content block as subtitle
        if slide.content_blocks:
            subtitle_placeholder = ppt_slide.placeholders[placeholders["subtitle"]]
            subtitle_placeholder.text = slide.content_blocks[0].text
            apply_text_styling(subtitle_placeholder.text_frame, theme, is_title=False)

    # Add speaker notes
    if slide.speaker_notes:
        notes_slide = ppt_slide.notes_slide
        notes_slide.notes_text_frame.text = slide.speaker_notes


def _add_content(ppt_slide, slide: Slide, content_idx: int, theme: Theme) -> None:
    """Add content to a slide.

    Args:
        ppt_slide: PowerPoint slide object
        slide: Slide with content
        content_idx: Index of content placeholder
        theme: Theme configuration
    """
    if not slide.content_blocks:
        return

    # Get content placeholder
    content_placeholder = ppt_slide.placeholders[content_idx]
    text_frame = content_placeholder.text_frame

    # Clear default content
    text_frame.clear()

    # Add content blocks
    for idx, block in enumerate(slide.content_blocks):
        # Add paragraph
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Set text
        p.text = block.text

        # Apply bullet formatting
        if block.is_bullet:
            p.level = min(block.indent_level, 4)  # PowerPoint supports up to 5 levels

        # Apply styling
        p.font.name = theme.font_body
        p.font.size = Pt(theme.font_size_body)
        p.font.color.rgb = hex_to_rgb(theme.text_color)
