"""Tests for deck builder module."""

import tempfile
from pathlib import Path

from pptx import Presentation

from scenarios.deck_builder import generate_presentation, parse_markdown


def test_parse_markdown_basic():
    """Test basic markdown parsing."""
    markdown = """## Slide 1: Title

This is content
- Bullet 1
- Bullet 2

**Speaker Notes:**
These are speaker notes

---

## Slide 2: Another Slide

More content here"""

    slides = parse_markdown(markdown)
    assert len(slides) == 2
    assert slides[0].title == "Title"
    assert slides[1].title == "Another Slide"
    assert slides[0].speaker_notes == "These are speaker notes"


def test_parse_markdown_with_bullets():
    """Test parsing markdown with bullet points."""
    markdown = """## Test Slide

- First bullet
  - Nested bullet
- Second bullet

**Speaker Notes:**
Notes here"""

    slides = parse_markdown(markdown)
    assert len(slides) == 1
    assert len(slides[0].content_blocks) == 3
    assert slides[0].content_blocks[0].is_bullet
    assert slides[0].content_blocks[0].text == "First bullet"
    assert slides[0].content_blocks[1].indent_level == 1


def test_generate_presentation():
    """Test generating a PowerPoint presentation."""
    markdown = """## Title Slide

Subtitle content

---

## Content Slide

- Point 1
- Point 2
- Point 3

**Speaker Notes:**
These are important notes"""

    slides = parse_markdown(markdown)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        output_path = Path(tmp.name)

    try:
        # Generate presentation
        generate_presentation(slides, output_path)

        # Verify file exists
        assert output_path.exists()

        # Verify presentation structure
        prs = Presentation(str(output_path))
        assert len(prs.slides) == 2

        # Check first slide
        slide1 = prs.slides[0]
        assert "Title Slide" in slide1.shapes.title.text if slide1.shapes.title else False

        # Check speaker notes
        slide2 = prs.slides[1]
        notes_frame = slide2.notes_slide.notes_text_frame
        if notes_frame:
            notes = notes_frame.text
            assert "These are important notes" in notes

    finally:
        # Clean up
        if output_path.exists():
            output_path.unlink()


def test_amplifier_walking_deck():
    """Test with the actual AMPLIFIER_WALKING_DECK.md file."""
    deck_path = Path("/workspaces/amplifier/docs/AMPLIFIER_WALKING_DECK.md")

    if not deck_path.exists():
        # Skip if file doesn't exist
        return

    with open(deck_path, "r", encoding="utf-8") as f:
        markdown_content = f.read()

    slides = parse_markdown(markdown_content)

    # Check we got the expected number of slides
    assert len(slides) > 10, f"Expected more than 10 slides, got {len(slides)}"

    # Check first slide
    assert "Amplifier" in slides[0].title or "Core Problem" in slides[0].title

    # Check speaker notes are preserved
    slides_with_notes = [s for s in slides if s.speaker_notes]
    assert len(slides_with_notes) > 5, "Expected multiple slides with speaker notes"

    # Test generation with actual deck
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        output_path = Path(tmp.name)

    try:
        generate_presentation(slides, output_path)
        assert output_path.exists()

        # Verify it's a valid PowerPoint file
        prs = Presentation(str(output_path))
        assert len(prs.slides) == len(slides)

    finally:
        if output_path.exists():
            output_path.unlink()
