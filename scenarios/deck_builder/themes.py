"""Theme system for presentations."""

from pptx.dml.color import RGBColor
from pptx.util import Pt

from .models import Theme


def apply_theme(presentation, theme: Theme):
    """Apply a theme to the presentation.

    Args:
        presentation: python-pptx Presentation object
        theme: Theme configuration
    """
    # python-pptx doesn't support global themes well
    # Styling is applied per slide/element in the generator
    # This function is kept for future enhancement
    # Currently theme is applied element-by-element in generator.py


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor.

    Args:
        hex_color: Hex color string (e.g., "#2C3E50")

    Returns:
        RGBColor object
    """
    # Remove # if present
    hex_color = hex_color.lstrip("#")
    # Convert to RGB
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def apply_text_styling(text_frame, theme: Theme, is_title: bool = False):
    """Apply theme styling to a text frame.

    Args:
        text_frame: python-pptx text frame
        theme: Theme configuration
        is_title: Whether this is title text
    """
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = theme.font_title if is_title else theme.font_body
            run.font.size = Pt(theme.font_size_title if is_title else theme.font_size_body)
            run.font.color.rgb = hex_to_rgb(theme.text_color)


def get_default_theme() -> Theme:
    """Get the default theme.

    Returns:
        Default Theme object
    """
    return Theme()
