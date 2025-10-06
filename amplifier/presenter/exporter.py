"""PowerPoint export module for creating .pptx files from presentations."""

from contextlib import suppress
from pathlib import Path

from pptx import Presentation as PptxPresentation
from presenter.models import Presentation
from presenter.models import Slide
from presenter.models import SlideType


class PowerPointExporter:
    """Exports Presentation objects to PowerPoint files."""

    def __init__(self):
        """Initialize the exporter."""
        self.prs = None
        self.slide_layouts = None

    def export(self, presentation: Presentation, output_path: str | Path) -> Path:
        """Export a presentation to a PowerPoint file.

        Args:
            presentation: Presentation object to export
            output_path: Output file path (.pptx)

        Returns:
            Path to the created file
        """
        output_path = Path(output_path)
        if not output_path.suffix:
            output_path = output_path.with_suffix(".pptx")

        # Create new PowerPoint presentation
        self.prs = PptxPresentation()
        self.slide_layouts = self._get_layouts()

        # Process each slide
        for slide in presentation.slides:
            self._add_slide(slide)

        # Save the presentation
        self.prs.save(output_path)
        return output_path

    def _get_layouts(self) -> dict:
        """Get available slide layouts."""
        return {
            "title": self.prs.slide_layouts[0],  # Title Slide
            "bullet": self.prs.slide_layouts[1],  # Title and Content
            "section": self.prs.slide_layouts[2],  # Section Header
            "two_content": self.prs.slide_layouts[3],  # Two Content
            "blank": self.prs.slide_layouts[5],  # Blank
            "content": self.prs.slide_layouts[1],  # Title and Content (fallback)
        }

    def _add_slide(self, slide: Slide):
        """Add a slide to the PowerPoint presentation."""
        if slide.type == SlideType.TITLE:
            self._add_title_slide(slide)
        elif slide.type == SlideType.SECTION:
            self._add_section_slide(slide)
        elif slide.type == SlideType.COMPARISON:
            self._add_comparison_slide(slide)
        elif slide.type == SlideType.CONCLUSION:
            self._add_conclusion_slide(slide)
        else:
            # Default to bullet slide
            self._add_bullet_slide(slide)

    def _add_title_slide(self, slide: Slide):
        """Add a title slide."""
        ppt_slide = self.prs.slides.add_slide(self.slide_layouts["title"])

        # Set title
        title_shape = ppt_slide.shapes.title
        if title_shape:
            title_shape.text = slide.title

        # Set subtitle if available
        if slide.content.main and len(slide.content.main) > 1:
            subtitle_shape = ppt_slide.placeholders[1] if len(ppt_slide.placeholders) > 1 else None
            if subtitle_shape:
                # Join additional content as subtitle
                subtitle_text = "\n".join(str(item) for item in slide.content.main[1:])
                subtitle_shape.text = subtitle_text

        # Add speaker notes
        if slide.content.notes:
            ppt_slide.notes_slide.notes_text_frame.text = slide.content.notes

    def _add_section_slide(self, slide: Slide):
        """Add a section divider slide."""
        ppt_slide = self.prs.slides.add_slide(self.slide_layouts["section"])

        # Set section title
        title_shape = ppt_slide.shapes.title
        if title_shape:
            title_shape.text = slide.title

        # Add speaker notes
        if slide.content.notes:
            ppt_slide.notes_slide.notes_text_frame.text = slide.content.notes

    def _add_bullet_slide(self, slide: Slide):
        """Add a standard bullet point slide."""
        ppt_slide = self.prs.slides.add_slide(self.slide_layouts["bullet"])

        # Set title
        title_shape = ppt_slide.shapes.title
        if title_shape:
            title_shape.text = slide.title

        # Add bullet points
        if slide.content.bullets:
            content_shape = None
            for shape in ppt_slide.shapes:
                if hasattr(shape, "text_frame") and shape != title_shape:
                    content_shape = shape
                    break

            if content_shape:
                text_frame = content_shape.text_frame
                text_frame.clear()  # Clear placeholder text

                for bullet in slide.content.bullets:
                    p = text_frame.add_paragraph()
                    p.text = bullet.strip()

                    # Handle indentation for sub-bullets
                    if bullet.startswith("  "):
                        p.level = 1
                        p.text = bullet.strip()[1:].strip()  # Remove bullet marker

                # Remove empty first paragraph
                if text_frame.paragraphs[0].text == "":
                    text_frame.paragraphs[0].text = " "
                    text_frame.paragraphs[0]._element.getparent().remove(text_frame.paragraphs[0]._element)

        # Add speaker notes
        if slide.content.notes:
            ppt_slide.notes_slide.notes_text_frame.text = slide.content.notes

    def _add_comparison_slide(self, slide: Slide):
        """Add a comparison slide with two columns."""
        ppt_slide = self.prs.slides.add_slide(self.slide_layouts["two_content"])

        # Set title
        title_shape = ppt_slide.shapes.title
        if title_shape:
            title_shape.text = slide.title

        # Add content to two columns
        content_shapes = [s for s in ppt_slide.shapes if hasattr(s, "text_frame") and s != title_shape]

        if slide.content.main and len(slide.content.main) >= 2:
            # Left column
            if len(content_shapes) > 0:
                self._add_bullets_to_shape(content_shapes[0], slide.content.main[0])

            # Right column
            if len(content_shapes) > 1:
                self._add_bullets_to_shape(content_shapes[1], slide.content.main[1])
        elif slide.content.bullets:
            # Split bullets into two columns
            mid = len(slide.content.bullets) // 2
            if len(content_shapes) > 0:
                self._add_bullets_to_shape(content_shapes[0], slide.content.bullets[:mid])
            if len(content_shapes) > 1:
                self._add_bullets_to_shape(content_shapes[1], slide.content.bullets[mid:])

        # Add speaker notes
        if slide.content.notes:
            ppt_slide.notes_slide.notes_text_frame.text = slide.content.notes

    def _add_conclusion_slide(self, slide: Slide):
        """Add a conclusion slide."""
        # Use bullet slide layout for conclusions
        self._add_bullet_slide(slide)

    def _add_bullets_to_shape(self, shape, bullets):
        """Add bullet points to a shape's text frame."""
        if not hasattr(shape, "text_frame"):
            return

        text_frame = shape.text_frame
        text_frame.clear()

        if isinstance(bullets, list):
            for bullet in bullets:
                p = text_frame.add_paragraph()
                p.text = str(bullet).strip()
                p.level = 0
        else:
            # Single item
            p = text_frame.add_paragraph()
            p.text = str(bullets)
            p.level = 0

        # Remove empty first paragraph if exists
        if text_frame.paragraphs and text_frame.paragraphs[0].text == "":
            with suppress(Exception):
                text_frame.paragraphs[0]._element.getparent().remove(text_frame.paragraphs[0]._element)
